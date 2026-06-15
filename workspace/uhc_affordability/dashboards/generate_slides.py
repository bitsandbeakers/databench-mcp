"""
Generate the UHC Affordability deck (.pptx) from spine v4 — BLUF-ordered, Optum-branded.

Audience = NON-TECHNICAL (the grading lens). So slide bodies are plain English; all the
technical depth (DRG, Spearman, Louvain, SHAP, falsification, CV/MAE) lives in the SPEAKER
NOTES, ready for a technical interviewer who probes. Figures come from app.py (single source
of truth) rendered via kaleido. Optum brand: Orange as fills only, Warm-Gray text.

Run: uv run --with python-pptx,kaleido,pyarrow python generate_slides.py
Out: UHC_Affordability_Deck.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import app  # figure functions + data

import sys
HERE = Path(__file__).parent
ASSETS = HERE / "slide_assets"; ASSETS.mkdir(exist_ok=True)
NO_NOTES = "--no-notes" in sys.argv  # read-ahead build: strip speaker notes
OUT = HERE / ("UHC_Affordability_Deck_readahead.pptx" if NO_NOTES else "UHC_Affordability_Deck.pptx")

ORANGE = RGBColor(0xFF, 0x61, 0x2B)   # fills/pops only, never type
INK    = RGBColor(0x3D, 0x3C, 0x38)   # all text (Warm Gray)
WWHITE = RGBColor(0xFA, 0xF8, 0xF2)
HORIZON= RGBColor(0xED, 0xE8, 0xE0)
FONT = "Arial"  # Enterprise Sans not bundled -> brand-approved fallback


def render(name, fig, w=900, h=560):
    p = ASSETS / f"{name}.png"
    try:
        fig.write_image(str(p), width=w, height=h, scale=2)
        return str(p)
    except Exception as e:
        print(f"  [skip image {name}: {e}]")
        return None


print("rendering figures...")
IMG = {
    "dispersion": render("dispersion", app.fig_dispersion()),
    "falsify":    render("falsify", app.fig_falsification()),
    "states":     render("states", app.fig_states()),
    "network":    render("network", app.build_network_fig(), w=1000, h=620),
    "chains":     render("chains", app.fig_chains()),
    "steer":      render("steer", app.fig_steer(), w=1000, h=560),
    "quality":    render("quality", app.fig_quality_scatter(), w=1000, h=560),
    "drivers":    render("drivers", app.fig_drivers()),
}

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _txt(tf, runs, size, bold=False, color=INK, align=PP_ALIGN.LEFT, space_after=6):
    tf.word_wrap = True
    if isinstance(runs, str):
        runs = [(runs, bold)]
    p = tf.paragraphs[0]; p.alignment = align; p.space_after = Pt(space_after)
    for t, b in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = b; r.font.name = FONT; r.font.color.rgb = color
    return p


def bg(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = WWHITE


def kicker_bar(slide, label):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(0.42), Inches(11.5), Inches(0.4))
    _txt(tb.text_frame, label.upper(), 13, bold=True, color=INK)


def title(slide, text, top=0.85, size=29, width=12.1):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(width), Inches(1.1))
    _txt(tb.text_frame, text, size, bold=True, color=INK)


def bullets(slide, items, left=0.6, top=2.0, width=6.0, height=4.7, size=15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        runs = it if isinstance(it, list) else [(it, False)]
        d = p.add_run(); d.text = "■  "; d.font.size = Pt(11); d.font.name = FONT; d.font.color.rgb = ORANGE
        for t, b in runs:
            r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b
            r.font.name = FONT; r.font.color.rgb = INK
    return tb


def picture(slide, path, left, top, width):
    if path:
        slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))


def step_strip(slide, steps, top=6.0, height=0.72, left=0.7, total=11.95, gap=0.2):
    """Plain-language horizontal process flow: light boxes + orange chevrons between."""
    n = len(steps); bw = (total - gap * (n - 1)) / n; x = left
    for i, st in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(bw), Inches(height))
        box.fill.solid(); box.fill.fore_color.rgb = HORIZON
        box.line.color.rgb = HORIZON; box.line.width = Pt(0.5); box.shadow.inherit = False
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = st; r.font.size = Pt(11); r.font.bold = True
        r.font.name = FONT; r.font.color.rgb = INK
        if i < n - 1:
            ar = slide.shapes.add_textbox(Inches(x + bw - 0.04), Inches(top + 0.10), Inches(gap + 0.08), Inches(0.5))
            ap = ar.text_frame.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
            rr = ap.add_run(); rr.text = "›"; rr.font.size = Pt(18); rr.font.bold = True
            rr.font.name = FONT; rr.font.color.rgb = ORANGE
        x += bw + gap


def notes(slide, text):
    if NO_NOTES:
        return
    slide.notes_slide.notes_text_frame.text = text


def new(kicker=None):
    s = prs.slides.add_slide(BLANK); bg(s)
    if kicker:
        kicker_bar(s, kicker)
    return s


def evidence(kicker, ttl, img, items, note):
    s = new(kicker); title(s, ttl)
    bullets(s, items, left=0.6, top=2.0, width=6.0, size=15)
    picture(s, img, 6.9, 1.9, 6.0)
    notes(s, note)
    return s


# ---- 1 · TITLE ----
s = new()
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
_txt(tf, "Hospital Pricing & Network Affordability", 40, bold=True, color=INK)
p = tf.add_paragraph(); p.space_before = Pt(10)
r = p.add_run(); r.text = "What the same care costs across 3,236 hospitals — and what a health plan can do about it"
r.font.size = Pt(18); r.font.name = FONT; r.font.color.rgb = INK
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.55), Inches(3.2), Inches(0.10))
band.fill.solid(); band.fill.fore_color.rgb = ORANGE; band.line.fill.background()
sb = s.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.0))
_txt(sb.text_frame, "Medicare hospital cost data · an AI-enabled analysis", 13, color=INK)
notes(s, "One-line hook: the same care, wildly different price. I'll give you the bottom line first, then the evidence. "
         "Data: public CMS Medicare hospital cost files (inpatient + outpatient), plus CMS ownership, rural/urban, "
         "location, and quality data. Audience note: keeping slides plain; happy to go as deep as you want on method.")

# ---- 2 · BLUF ----
s = new("The bottom line")
title(s, "The same care can cost up to 16× more — and the gap is mostly pricing, not better care")
bullets(s, [
    [("Wildly different prices. ", True), ("The same hospital treatment can cost up to 16× more at one hospital "
     "than another.", False)],
    [("It's how they price, not what care costs. ", True), ("For-profit hospitals charge ~45% more — yet Medicare "
     "pays them slightly less. So it's a pricing choice, not higher costs.", False)],
    [("Paying more doesn't buy better care. ", True), ("Price and hospital quality ratings have no relationship — and "
     "the most expensive hospitals tend to rate the worst.", False)],
    [("The opportunity. ", True), ("Guide members to hospitals that are both cheaper and higher-quality "
     "(738 qualify).", False)],
    [("The size. ", True), ("On for-profit hospitals alone, the overcharge is roughly $6–26 billion.", False)],
], left=0.6, top=1.9, width=12.2, size=17)
rec = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.8))
rec.fill.solid(); rec.fill.fore_color.rgb = HORIZON; rec.line.fill.background()
rt = rec.text_frame; rt.word_wrap = True; rt.margin_left = Inches(0.2)
_txt(rt, [("RECOMMENDATION:  ", True), ("build the network around hospitals that are both affordable and high-quality — "
     "start where a cheaper, better option is already nearby.", False)], 15, color=INK)
notes(s, "If you remember one slide, this is it — four findings build to one action, impact up front. "
         "Technical backing for each, in order: (1) within-DRG charge dispersion, p90/p10 median 3.67×, max ~16× (DRG 470). "
         "(2) For-profit premium +45% on submitted charge vs −9% on wage-index-adjusted Medicare payment — a falsification "
         "test isolating pricing from cost. (3) Spearman(cost index, CMS star) = −0.03. (4) Quality-gated steerage. "
         "(5) $25.8B excess premium, $12.9B at 50% recovery; I show 25/50/100% scenarios.")

# ---- 3 · cost variation (Insight Q1) ----
evidence("Finding 1 · price varies wildly", "The same procedure, very different price",
    IMG["dispersion"], [
        [("The question: ", True), ("which treatments vary most in price between hospitals?", False)],
        "For the same treatment, the priciest hospitals charge at least 3× more than the cheapest — true for 94% of common treatments.",
        "The most extreme: a hip/knee replacement ran from about $20k to $316k — same surgery.",
        "It's the hospital you walk into, not the care, that drives the difference.",
    ],
    "Insight Discovery, question 1 (cost variation across providers). Method: within-DRG dispersion so the procedure is "
    "held constant; median p90/p10 = 3.67×, 252/269 DRGs ≥3×. Worked example DRG 470 (major joint replacement), "
    "$19,972 → $315,717. Charges are skew-12.5, so everything uses medians/logs, never raw means. This is 'submitted "
    "charges' (the hospital's list price) — see the data slide for why that's the right signal for a %-of-charges payer.")

# ---- 4 · pricing not cost ----
evidence("Why — the mechanism", "It's a pricing choice, not a cost difference",
    IMG["falsify"], [
        "For-profit hospitals charge about 45% more than others for comparable care — even after accounting for how sick "
        "patients are, where the hospital is, and what kind of hospital it is.",
        [("The tell: ", True), ("Medicare pays those same hospitals slightly LESS, not more.", False)],
        "Medicare's payment already adjusts for local labor costs — so 'things just cost more here' can't explain the gap.",
        "This only hurts an insurer that pays a share of the hospital's list price.",
    ],
    "The acid test (falsification). Under identical controls the for-profit gap is +45% on submitted charge but −9% on "
    "formula-set Medicare payment. If it were genuine cost, the cost-based (wage-index-adjusted) Medicare payment would "
    "rise, not fall — so cost-of-living can't be the driver. Nested OLS on log charge; premium 62.5%→46.3% as controls "
    "are added. Pre-empts the obvious 'urban hospitals cost more' objection.")

# ---- 5 · regional patterns (Insight Q2) ----
evidence("Finding 2 · geography", "Where you are matters — a lot",
    IMG["states"], [
        [("The question: ", True), ("do prices follow a regional pattern?", False)],
        "Yes — the most expensive states charge about 6× more than the cheapest for the same care.",
        "Maryland, which regulates hospital rates, is the cheapest; Nevada, California, New Jersey and Florida are the priciest.",
        "Same treatment, very different price depending on the market.",
    ],
    "Insight Discovery, question 2 (regional patterns). State case-mix-adjusted charge index 0.35 (MD) → 2.14 (NV), ~6×. "
    "Maryland's all-payer rate-setting waiver is the natural experiment — lowest charges, highest regulated payment. "
    "In the driver model geography is 18% of charge variance vs only 5.4% of cost-based payment variance.")

# ---- 6 · network (centerpiece) ----
evidence("Comparing fairly", "I group hospitals by the kind of care they actually do",
    IMG["network"], [
        "I sorted all 3,236 hospitals into families by the mix of care they deliver — academic medical centers, "
        "community hospitals, surgical specialty centers, rural hospitals, and so on.",
        "The groups came out of the data automatically — not hand-picked.",
        [("Why it matters: ", True), ("I compare each hospital only to its true peers, so 'our patients are just "
         "sicker' doesn't hold up.", False)],
        "In the live dashboard you can pick a hospital type and see the flagged ones.",
    ],
    "The network-analytics centerpiece (role fit) — linger and demo the live dropdowns. Method: cosine similarity of "
    "each hospital's service-mix vector → k-nearest-neighbour graph → Louvain community detection, 14 communities at "
    "modularity 0.70. Archetype-relative comparison is what makes the outlier list defensible (drops Cedars-Sinai as "
    "'expensive like its academic peers', keeps Stanford).")

# ---- 7 · outliers both IP & OP ----
evidence("The overchargers", "Which hospitals charge far more than their peers",
    IMG["chains"], [
        "I flagged hospitals that charge far above their true peers — in both hospital stays and outpatient care.",
        "The same names show up in both: a handful of for-profit chains, plus a few high-end academic centers.",
        [("They cluster: ", True), ("a couple of for-profit chains are flagged ~100% of the time; non-profits almost "
         "never are.", False)],
        "It's a whole-hospital habit — flagged hospitals overcharge across nearly all their services, not just a few.",
    ],
    "Task instruction 3 (outlier detection), explicitly in BOTH settings + patterns. Inpatient = within-DRG top-charge "
    "providers (h006: Stanford, Cedars, HCA); outpatient = within-APC (h007: same names). Patterns (h008): cluster by "
    "state (CA/NJ/FL) and for-profit chain — Capital Health & Carepoint 100%, HCA 24% vs independents 4.5%. Repeat-"
    "offenders price ~95% of their service lines above the national 90th percentile.")

# ---- 8 · steerage ----
evidence("The action", "Guide members to a cheaper, comparable hospital nearby",
    IMG["steer"], [
        "'Steering' means pointing members to a cheaper, comparable hospital — and it has to be realistically close.",
        "For 9 of the 11 worst offenders, there's a cheaper, comparable hospital within ~40 miles (one saves 87%).",
        "For 2, there's no nearby alternative — steering can't help there; you negotiate rates instead.",
    ],
    "Steerage is TOWARD a lower-cost comparable peer (corrected from 'away from'), and it's inherently local "
    "(40-mi great-circle from ZIP centroids as a drive-time proxy). 9/11 locally steerable; the 2 misses are a "
    "network-adequacy gap (199–345 mi), an honest finding. Production needs true drive-time + in-network + adequacy.")

# ---- 9 · quality gate ----
evidence("The quality check", "Paying more doesn't buy better care",
    IMG["quality"], [
        "I added the government's hospital quality star ratings (public data).",
        [("Price and quality are unrelated ", True), ("— and the most expensive hospitals actually rate the worst.", False)],
        "Most of our flagged hospitals are both expensive AND low-rated — steering away wins twice.",
        "Two are expensive but genuinely excellent (Stanford, MLK) — the quality check keeps them off the avoid list.",
    ],
    "The strongest finding — and it was in-scope public CMS data (Hospital Compare). Spearman(cost, star) = −0.03; the "
    "≥2× cost tier averages 2.74★ vs ~3.2 elsewhere. 7 of 10 rated flagged hospitals are 1–2★. The gate correctly "
    "protects the high-quality expensive ones. Cost is necessary, not sufficient — this closes the 'cost-only steerage "
    "is unsafe' objection.")

# ---- 10 · drivers (bonus) ----
evidence("Bonus · what drives price", "Location, type of care, and ownership — not size",
    IMG["drivers"], [
        "What actually predicts a hospital's price: where it is, the type and complexity of care, and who owns it.",
        "Hospital size doesn't matter — bigger isn't cheaper per service.",
        "Three independent methods all agree on the top drivers, so we trust the ranking.",
        "The model predicts a hospital's cost to within about $11k (roughly 22%).",
    ],
    "Bonus question (cost predictors). Triangulated importance — interventional SHAP + permutation + glass-box EBM — "
    "perfectly seed-stable: state > case-mix > ownership; volume negligible. Reported with 5-fold CV MAE + dollar error, "
    "not R² alone (R² on a log target flatters fit). EBM vs LightGBM overlap within ~1 SD. All associational.")

# ---- 11 · approach, data & rigor ----
s = new("How it was done"); title(s, "Two datasets, combined by the question — then supplemented")
bullets(s, [
    [("The core (the assignment): ", True), ("two public datasets — every hospital's inpatient (hospital-stay) and "
     "outpatient costs. I used both, throughout.", False)],
    [("Explored & cleaned first: ", True), ("I profiled every table up front — distributions, missing values, outliers — "
     "fixed file-format/encoding issues, dropped CMS-suppressed and unreliable rows, and caught a data artifact before "
     "it could mislead.", False)],
    [("Combined several ways, depending on the question: ", True), ("side-by-side per hospital (are hospitals pricey "
     "for stays also pricey outpatient? yes), pooled together for overall patterns, and blended into each hospital's "
     "'care fingerprint' that defines the hospital families.", False)],
    [("Then supplemented ", True), ("with more government data to answer the harder, more interesting questions: "
     "ownership, rural-vs-urban, location/distance, and quality ratings.", False)],
    [("Assumption + check: ", True), ("list price as the affordability signal; and every headline finding repeats in "
     "the prior year — real patterns, not a one-year fluke.", False)],
], left=0.6, top=1.9, width=12.2, size=15)
notes(s, "Task instruction 1 (data understanding & prep / INTEGRATION). The core ask was the inpatient + outpatient "
         "files — I used both, and combined them FOUR ways by analysis: (1) separately at natural grain (within-DRG "
         "inpatient + within-APC outpatient) for outlier detection; (2) stacked long (~210k rows) for pooled cost-index "
         "and markup; (3) merged per-hospital, both settings, for the pricing-culture correlation (IP vs OP charge index "
         "ρ=0.82) and the driver model; (4) a combined DRG+APC service-mix matrix for the archetype network. THEN "
         "supplemented with CMS HCRIS (ownership, rural margins → for-profit premium + rural exposure), RUCA (urbanicity), "
         "ZIP centroids (distance → in-market steerage), and Hospital Compare (quality). Cleaning: Latin-1→UTF-8, DECIMAL "
         "coercion, CCN key, RUCA=99 drop, cell-suppression noted, HCRIS margin fix. Stack: DuckDB, pandas, sklearn, "
         "statsmodels, plotly. Robustness = 2022 out-of-year holdout.")

# ---- 12 · AI enablement ----
s = new("How it was built"); title(s, "AI-enabled analysis — a custom MCP server + an analyst skill")
bullets(s, [
    [("I built an MCP server (databench-mcp): ", True), ("a guard-railed tool surface the AI operates — ingest, profile, "
     "explore, model, network-analyze, visualize — with discipline built in (profile before modeling, no leakage).", False)],
    [("An analyst skill defined and required the process below: ", True), ("nothing skipped or hand-waved; I set the "
     "hypotheses and corrected the approach when it was wrong.", False)],
    [("Adversarially verified: ", True), ("findings stay 'suspect' until they survive a refutation check — this caught a "
     "self-prediction (data-leakage) bug and a data artifact before either reached a conclusion.", False)],
    [("Evidence-tracked & self-improving: ", True), ("18 of 18 conclusions trace to a logged result; the corrections "
     "become the tool's own to-do list.", False)],
], left=0.6, top=1.8, width=12.2, size=14)
lbl = s.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12.2), Inches(0.35))
_txt(lbl.text_frame, "The process the skill enforced, every time:", 12, bold=True, color=INK)
step_strip(s, ["Understand", "Get & clean data", "Explore", "Test each idea", "Verify", "Recommend"], top=5.92)
notes(s, "Top line stays 'How it was built'; the differentiator is AI enablement. The step strip is plain-language but "
         "maps to CRISP-DM (Business + Data Understanding → Data Prep → Modeling → Evaluation → Deployment), enforced by "
         "the databench-analyst skill. Adversarial verification = findings default to suspect until they survive "
         "refutation (caught the charge-as-its-own-predictor leakage and the HCRIS total-margin artifact). 13-entry "
         "correction ledger; product-gap corrections become the tool's roadmap. If a technical interviewer wants the "
         "framework, CRISP-DM is the reference — happy to show the phase mapping.")

# ---- 13 · limits & next ----
s = new("Future improvements"); title(s, "Where this goes next — with real plan data")
hA = s.shapes.add_textbox(Inches(0.6), Inches(1.65), Inches(6.2), Inches(0.4))
_txt(hA.text_frame, "Analysis & data", 16, bold=True, color=INK)
bullets(s, [
    [("Care quality: ", True), ("do the factors that raise price also affect quality? Model both together.", False)],
    [("Prove cause, ", True), ("not just correlation.", False)],
    [("Network-aware steering: ", True), ("reachable, in-network, right service line, has capacity.", False)],
    [("Real prices: ", True), ("negotiated insurer rates and member claims (we used Medicare list prices).", False)],
    [("Access & equity: ", True), ("don't steer a community away from its only nearby hospital.", False)],
    [("Automate & refresh: ", True), ("rebuild each year on fresher location and ownership data.", False)],
], left=0.6, top=2.15, width=6.2, size=13)
hB = s.shapes.add_textbox(Inches(7.05), Inches(1.65), Inches(5.7), Inches(0.4))
_txt(hB.text_frame, "Strengthen the AI platform", 16, bold=True, color=INK)
bullets(s, [
    [("Expand the model set ", True), ("— plus a skill that picks the right method(s) for each question.", False)],
    [("Prevent repeat mistakes: ", True), ("turn each correction into a guardrail.", False)],
    [("Privacy & compliance guardrails ", True), ("(HIPAA / PHI, GDPR) built into the platform.", False)],
    [("A shared knowledge base ", True), ("(wiki) so the method and evidence are reusable.", False)],
], left=7.05, top=2.15, width=5.7, size=13)
notes(s, "Reorganized: Analysis & data vs the AI platform (enablement) — model-selection, prevention, and the wiki are "
         "all enablement. Technical detail: quality = drivers-of-quality + a joint multivariate price×quality model "
         "(cost–quality efficiency frontier); cause = Double-ML / causal forest (HHI + wage-index confounders); network-"
         "aware = drive-time + CMS time/distance adequacy + in-network + service-line + capacity; real prices = "
         "Transparency-in-Coverage MRFs + plan claims; access & equity = before down-tiering a community's only nearby "
         "hospital (the safety-net 'no local option' cases, e.g. MLK Jr), check the access impact; productionize = "
         "re-runnable recipe + drift monitoring, current ZIP/ZCTA + chain crosswalk + Critical Access Hospitals. AI "
         "platform: an expanded model library + a method-selection skill; a 'prevention' step that turns each correction "
         "into an automated guardrail (not just a tool fix); a knowledge-base wiki as institutional memory.")

# ---- 14 · recommendation ----
s = new("Recommendation"); title(s, "What we can do")
bullets(s, [
    [("Build the network on cost AND quality, not cost alone", True)],
    [("Start where a cheaper, better hospital already exists nearby", True)],
    [("Where there's no nearby option, negotiate rates instead", True)],
], left=0.6, top=2.3, width=12.2, size=20)
lb = s.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.6))
_txt(lb.text_frame, [("Interactive leave-behind:  ", True),
     ("bitsandbeakers.github.io/databench-mcp/uhc-affordability/", False)], 14, color=INK)
notes(s, "Voice track (bullets are headlines only): (1) Cost-only steering can send members to cheap-but-bad hospitals — "
         "the quality gate prevents that. (2) Start with the overpriced, low-quality hospitals that already have a "
         "cheaper, higher-quality neighbor within ~40 miles — the fast wins. (3) For the markets with no nearby "
         "alternative (the 2 'no local option' cases), steering can't help — negotiate rates or redesign the network "
         "there. Close on action, not findings; point them to the live dashboard as the interactive leave-behind.")

prs.save(str(OUT))
print(f"\nwrote {OUT.name} — {len(prs.slides._sldIdLst)} slides")
